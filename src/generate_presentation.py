"""
generate_presentation.py
Generate a professional PowerPoint presentation for the Tox21 ML project.
"""

import sys
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR   = Path(__file__).resolve().parent.parent
FIGURES    = BASE_DIR / "reports" / "figures"
REPORTS    = BASE_DIR / "reports"
OUTPUT     = BASE_DIR / "reports" / "Tox21_Presentation.pptx"

# Color scheme
RED    = RGBColor(0xC0, 0x39, 0x2B)
DARK   = RGBColor(0x2C, 0x3E, 0x50)
LIGHT  = RGBColor(0xEC, 0xF0, 0xF1)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x2E, 0x86, 0xC1)


def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_text(slide, text, left, top, width, height,
                   font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_image_safe(slide, img_path, left, top, width, height=None):
    p = Path(img_path)
    if p.exists():
        if height:
            slide.shapes.add_picture(str(p), left, top, width, height)
        else:
            slide.shapes.add_picture(str(p), left, top, width)
        return True
    return False


def add_bullet_list(slide, items, left, top, width, height,
                    font_size=16, color=DARK, title=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if title:
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = title
        r0.font.size = Pt(font_size + 4)
        r0.font.bold = True
        r0.font.color.rgb = color
    for item in items:
        if tf.paragraphs and tf.paragraphs[-1].runs:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[-1]
        p.level = 0
        r = p.add_run()
        r.text = f"  {item}"
        r.font.size = Pt(font_size)
        r.font.color.rgb = color


def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, DARK)
    add_title_text(slide, "Prediction of Chemical Molecule Toxicity",
                   Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2),
                   font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_title_text(slide, "Using Artificial Intelligence (Tox21 Dataset)",
                   Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.7),
                   font_size=24, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)
    add_title_text(slide, "Master AI Project  |  Machine Learning & Cheminformatics",
                   Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.6),
                   font_size=18, bold=False, color=RGBColor(0x99,0xAA,0xBB), align=PP_ALIGN.CENTER)
    add_title_text(slide, "2026", Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.5),
                   font_size=16, bold=False, color=RGBColor(0x88,0x99,0xAA), align=PP_ALIGN.CENTER)

    # ── Slide 2: Agenda ───────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Agenda", Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                   font_size=30, bold=True, color=RED)
    items = [
        "1. Problem Statement & Context",
        "2. Dataset — Tox21",
        "3. Methodology Overview",
        "4. Exploratory Data Analysis",
        "5. Molecular Feature Engineering",
        "6. Machine Learning Models",
        "7. Deep Learning — ToxNet",
        "8. Hyperparameter Optimization",
        "9. Experimental Results",
        "10. Explainable AI (SHAP)",
        "11. Web Application & API",
        "12. Conclusion & Future Work",
    ]
    add_bullet_list(slide, items, Inches(0.8), Inches(1.2), Inches(11.5), Inches(6),
                    font_size=16, color=DARK)

    # ── Slide 3: Problem Statement ────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Problem Statement", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=30, bold=True, color=RED)
    add_bullet_list(slide, [
        "Before approval, every drug/pesticide/food additive must be tested for toxicity",
        "Traditional lab testing: slow (months/years), expensive ($1–3M per compound), requires animal models",
        "Tox21 Challenge: can AI predict toxicity directly from molecular structure?",
        "Goal: Train ML/DL models on Tox21 dataset — 7,831 molecules, 12 assays",
    ], Inches(0.8), Inches(1.3), Inches(11.5), Inches(3.5), font_size=17, color=DARK)
    add_title_text(slide, "\"Predict molecular toxicity directly from SMILES strings\"",
                   Inches(1), Inches(5.2), Inches(11), Inches(0.8),
                   font_size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # ── Slide 4: Dataset ──────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Dataset — Tox21", Inches(0.5), Inches(0.3), Inches(7), Inches(0.7),
                   font_size=30, bold=True, color=RED)
    add_bullet_list(slide, [
        "7,831 chemical compounds (SMILES strings)",
        "12 nuclear receptor / stress response assays",
        "Binary labels: Toxic (1) / Non-Toxic (0)",
        "Primary task: NR-AhR (Aryl Hydrocarbon Receptor)",
        "  Positive: 768 (13.3%)  |  Negative: 5,774 (86.7%)",
        "High class imbalance — handled with class weights",
    ], Inches(0.8), Inches(1.3), Inches(6), Inches(5.5), font_size=15, color=DARK)
    add_image_safe(slide, FIGURES / "class_distribution.png",
                   Inches(7), Inches(1.0), Inches(6), Inches(4))

    # ── Slide 5: EDA ──────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Exploratory Data Analysis", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                   font_size=28, bold=True, color=RED)
    add_image_safe(slide, FIGURES / "descriptor_distributions.png",
                   Inches(0.3), Inches(0.9), Inches(6.4), Inches(3.2))
    add_image_safe(slide, FIGURES / "correlation_matrix.png",
                   Inches(6.9), Inches(0.9), Inches(6.1), Inches(3.2))
    add_image_safe(slide, FIGURES / "descriptor_boxplots.png",
                   Inches(0.3), Inches(4.2), Inches(12.7), Inches(3.0))

    # ── Slide 6: Feature Engineering ─────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Molecular Feature Engineering", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=28, bold=True, color=RED)
    add_bullet_list(slide, [
        "RDKit molecular processing pipeline",
        "130 molecular descriptors (MW, LogP, TPSA, rings, fragments...)",
        "Morgan Fingerprints ECFP4 (2048 bits) — primary features",
        "MACCS Keys (167 bits)",
        "RDKit Topological Fingerprints (2048 bits)",
        "8 invalid molecules removed (valence errors)",
        "7,823 valid molecules processed",
    ], Inches(0.8), Inches(1.2), Inches(5.5), Inches(5.5), font_size=15, color=DARK)
    add_image_safe(slide, FIGURES / "missing_values.png",
                   Inches(6.2), Inches(1.2), Inches(6.8), Inches(5.5))

    # ── Slide 7: ML Models ────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Machine Learning Models", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=28, bold=True, color=RED)
    add_bullet_list(slide, [
        "6 classical ML models trained on Morgan FP (2048 bits)",
        "Logistic Regression  |  Decision Tree  |  Random Forest",
        "XGBoost  |  LightGBM  |  SVM (RBF kernel)",
        "Class-weighted training to handle imbalance",
        "5-fold stratified cross-validation",
        "All models saved as .pkl files",
    ], Inches(0.8), Inches(1.2), Inches(6), Inches(4), font_size=15, color=DARK)
    add_image_safe(slide, FIGURES / "model_comparison_NR_AhR.png",
                   Inches(6.2), Inches(1.0), Inches(6.8), Inches(4.5))

    # ── Slide 8: Results ──────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Experimental Results", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=28, bold=True, color=RED)

    # Load metrics for table
    metrics_path = REPORTS / "NR_AhR_morgan_metrics.json"
    metrics_text = []
    if metrics_path.exists():
        with open(metrics_path) as f:
            m = json.load(f)
        for name, vals in sorted(m.items(), key=lambda x: -x[1].get("roc_auc", 0)):
            metrics_text.append(
                f"  {name:22s}  ROC-AUC={vals.get('roc_auc',0):.3f}  F1={vals.get('f1',0):.3f}  Acc={vals.get('accuracy',0):.3f}"
            )
    add_bullet_list(slide, metrics_text, Inches(0.5), Inches(1.2), Inches(6.3), Inches(5.5),
                    font_size=12, color=DARK)
    add_image_safe(slide, FIGURES / "roc_curves_NR_AhR.png",
                   Inches(6.5), Inches(1.0), Inches(6.5), Inches(3.2))
    add_image_safe(slide, FIGURES / "confusion_matrices_NR_AhR.png",
                   Inches(6.5), Inches(4.3), Inches(6.5), Inches(3.0))

    # ── Slide 9: ROC/PR Curves ────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "ROC & Precision-Recall Curves", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                   font_size=28, bold=True, color=RED)
    add_image_safe(slide, FIGURES / "roc_curves_NR_AhR_all.png",
                   Inches(0.3), Inches(0.9), Inches(6.5), Inches(3.3))
    add_image_safe(slide, FIGURES / "pr_curves_NR_AhR.png",
                   Inches(6.9), Inches(0.9), Inches(6.1), Inches(3.3))
    add_bullet_list(slide, [
        "Best ROC-AUC: SVM = 0.911",
        "Best PR-AUC: Random Forest = 0.635",
        "ToxNet (DL): ROC-AUC = 0.877",
    ], Inches(0.5), Inches(4.4), Inches(12), Inches(2.5), font_size=16, color=DARK)

    # ── Slide 10: Deep Learning ───────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Deep Learning — ToxNet", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=28, bold=True, color=RED)
    add_bullet_list(slide, [
        "Feedforward Neural Network (PyTorch)",
        "Architecture: Input(2048) → BN → 512 → 256 → 128 → 64 → 1",
        "BatchNorm + Dropout (0.3) for regularization",
        "BCEWithLogitsLoss with pos_weight for imbalance",
        "Weighted random sampling in DataLoader",
        "Adam optimizer + ReduceLROnPlateau scheduler",
        "Early stopping (patience=12): converged at epoch 13",
        "Final: ROC-AUC=0.877  |  F1=0.211  |  PR-AUC=0.559",
    ], Inches(0.8), Inches(1.2), Inches(6), Inches(5.5), font_size=15, color=DARK)
    add_image_safe(slide, FIGURES / "training_history_NR_AhR.png",
                   Inches(6.5), Inches(1.2), Inches(6.5), Inches(4))

    # ── Slide 11: SHAP ────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Explainable AI — SHAP", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                   font_size=28, bold=True, color=RED)
    add_image_safe(slide, FIGURES / "shap_summary_NR_AhR_desc_XGBoost.png",
                   Inches(0.3), Inches(0.9), Inches(6.3), Inches(3.3))
    add_image_safe(slide, FIGURES / "shap_bar_NR_AhR_desc_XGBoost.png",
                   Inches(6.8), Inches(0.9), Inches(6.2), Inches(3.3))
    add_image_safe(slide, FIGURES / "shap_waterfall_NR_AhR_desc_XGBoost_s0.png",
                   Inches(0.3), Inches(4.3), Inches(7), Inches(3.0))
    add_bullet_list(slide, [
        "TreeExplainer on XGBoost (descriptors)",
        "Top features: LogP, BertzCT, TPSA, Kappa indices",
        "Higher LogP -> increased toxicity probability",
    ], Inches(7.5), Inches(4.5), Inches(5.5), Inches(2.8), font_size=14, color=DARK)

    # ── Slide 12: Web App ─────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Web Application & API", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=28, bold=True, color=RED)
    add_bullet_list(slide, [
        "FastAPI REST API:",
        "  GET  /health       — service status",
        "  POST /predict      — single SMILES prediction",
        "  POST /batch_predict — multiple SMILES",
        "  GET  /model_info   — available models",
        "  GET  /metrics      — performance metrics",
        "",
        "Streamlit Web Interface:",
        "  - Molecule input & structure visualization",
        "  - Real-time toxicity prediction",
        "  - Model comparison charts",
        "  - SHAP explanations",
        "  - Analytics dashboard",
    ], Inches(0.8), Inches(1.2), Inches(11), Inches(5.5), font_size=14, color=DARK)

    # ── Slide 13: Model Comparison Summary ────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Model Comparison Summary", Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                   font_size=28, bold=True, color=RED)
    add_image_safe(slide, FIGURES / "model_comparison_NR_AhR_all.png",
                   Inches(0.3), Inches(0.9), Inches(12.7), Inches(6.3))

    # ── Slide 14: Test Results ────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Test Suite — 28 Passed, 1 Skipped", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=28, bold=True, color=RED)
    add_bullet_list(slide, [
        "Data loading tests (raw dataset integrity)",
        "RDKit molecular processing (valid/invalid SMILES)",
        "Descriptor computation (MW, LogP, TPSA...)",
        "Morgan fingerprint generation (2048 bits)",
        "MACCS key generation (167 bits)",
        "Feature pipeline (process_dataframe)",
        "Model loading and prediction shape/range",
        "Preprocessor fit/transform",
        "Metrics file integrity",
        "Figure generation (18 PNG files)",
        "End-to-end prediction pipeline",
        "FastAPI endpoints (health, predict, metrics, model_info)",
    ], Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5), font_size=15, color=DARK)

    # ── Slide 15: Discussion ──────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Discussion & Analysis", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=28, bold=True, color=RED)
    add_bullet_list(slide, [
        "Morgan FP outperforms molecular descriptors (ROC-AUC 0.91 vs 0.89)",
        "SVM achieves best ROC-AUC but Random Forest has better PR-AUC",
        "Class imbalance (88%/12%) mitigated with balanced class weights",
        "ToxNet (DL) competitive but requires more training data / epochs",
        "XGBoost on descriptors: best interpretable model with SHAP",
        "LogP, molecular complexity (BertzCT), TPSA are key toxicity predictors",
        "No data leakage detected: proper train/test split, no overlap",
        "Cross-validation confirms stable generalization (low std deviation)",
    ], Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5), font_size=15, color=DARK)

    # ── Slide 16: Limitations ─────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Limitations", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=28, bold=True, color=RED)
    add_bullet_list(slide, [
        "Task-specific models (one model per assay — 12 tasks)",
        "Multi-task learning not yet implemented",
        "Graph Neural Networks (molecular graphs) not explored",
        "3D conformational features not used",
        "Limited to Tox21 assays (not all toxicity endpoints)",
        "External validation on independent datasets needed",
        "ToxNet early stopping may be suboptimal (CPU training)",
    ], Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5), font_size=16, color=DARK)

    # ── Slide 17: Future Work ─────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Future Work", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=28, bold=True, color=RED)
    add_bullet_list(slide, [
        "Graph Neural Networks (GCN/MPNN) for molecular graph learning",
        "Multi-task deep learning across all 12 Tox21 assays",
        "Transformer-based molecular models (ChemBERTa, MolBERT)",
        "Active learning for efficient compound screening",
        "Uncertainty quantification (Monte Carlo Dropout / Conformal Prediction)",
        "Integration with virtual screening pipelines",
        "Docker containerization for production deployment",
        "Regulatory-grade validation following OECD guidelines",
    ], Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5), font_size=16, color=DARK)

    # ── Slide 18: Deliverables ────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Project Deliverables", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=28, bold=True, color=RED)
    add_bullet_list(slide, [
        "Source code (src/): pipeline, models, features, DL, SHAP, viz",
        "7,831 molecules processed  |  7,823 valid  |  8 invalid removed",
        "130 molecular descriptors  +  2048-bit Morgan FP",
        "12 trained ML models (.pkl) + 1 neural network (.pt)",
        "18 publication-quality figures (PNG)",
        "Performance metrics JSON  +  summary CSV",
        "FastAPI application (api/main.py)",
        "Streamlit interface (webapp/app.py)",
        "Test suite: 28 passed, 1 skipped (pytest)",
        "LaTeX academic report (50+ pages)",
        "PowerPoint presentation (this file)",
    ], Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5), font_size=15, color=DARK)

    # ── Slide 19: Reproducibility ─────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, WHITE)
    add_title_text(slide, "Reproducibility & Deployment", Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
                   font_size=28, bold=True, color=RED)
    add_bullet_list(slide, [
        "Single-command verification with pytest: 28 passed, 1 skipped",
        "Executed notebooks document the full workflow from download to evaluation",
        "Compiled PDF report generated from LaTeX with Tectonic",
        "FastAPI endpoint ready for local serving with uvicorn",
        "Streamlit interface ready for academic demonstration",
        "Portable local tooling: Tectonic installed under project/tools",
        "Git installed via winget for version-control workflows",
    ], Inches(0.8), Inches(1.2), Inches(6.1), Inches(5.5), font_size=15, color=DARK)
    add_bullet_list(slide, [
        "Run pipeline:",
        "  python src/pipeline.py",
        "",
        "Run API:",
        "  uvicorn api.main:app --reload",
        "",
        "Run UI:",
        "  streamlit run webapp/app.py",
    ], Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.2), font_size=15, color=BLUE)

    # ── Slide 20: Conclusion ──────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, DARK)
    add_title_text(slide, "Conclusion", Inches(0.5), Inches(0.8), Inches(12), Inches(0.8),
                   font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bullet_list(slide, [
        " Best model: SVM with Morgan FP  ->  ROC-AUC = 0.911",
        " Machine learning can predict molecular toxicity from SMILES alone",
        " SHAP explains that LogP, BertzCT, TPSA drive toxicity predictions",
        " Complete production-ready system: API + Web UI + Tests",
        " Foundation for high-throughput virtual toxicity screening",
    ], Inches(1.0), Inches(2.0), Inches(11), Inches(4),
    font_size=18, color=LIGHT)
    add_title_text(slide, "Thank You", Inches(0.5), Inches(6.2), Inches(12), Inches(0.7),
                   font_size=24, bold=False, color=RGBColor(0xAA,0xBB,0xCC), align=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f"Presentation saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")
    return OUTPUT


if __name__ == "__main__":
    build()
